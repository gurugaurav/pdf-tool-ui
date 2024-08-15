from docx import Document
import fitz
from PIL import Image
from PyPDF2 import PdfFileMerger, PdfFileReader, PdfFileWriter, PdfReader, PdfWriter
from fastapi import UploadFile
from fpdf import FPDF
import io
from typing import List  # Import List from typing

from pptx import Presentation
from pptx.util import Inches

PAGE_SIZES = {
    "A3": (297, 420),
    "A4": (210, 297),
    "A5": (148, 210),
    "Letter": (216, 279),
    "Legal": (216, 356),
}


def get_page_size(page_size):
    if isinstance(page_size, tuple):
        return page_size
    return PAGE_SIZES.get(page_size, PAGE_SIZES["A4"])


def convert_images_to_pdf(files: List[UploadFile], orientation='P', format='a4', margin=0):
    pdf = FPDF(orientation=orientation, unit="mm", format=format.lower())

    for file in files:
        image = Image.open(file.file)
        image_width_px, image_height_px = image.size

        # Convert pixel dimensions to mm (assuming 96 DPI)
        DPI = 96
        image_width_mm = image_width_px * 25.4 / DPI
        image_height_mm = image_height_px * 25.4 / DPI

        # Define page size in mm
        page_width_mm, page_height_mm = get_page_size(format.upper())

        # Calculate available space after applying margins
        available_width_mm = page_width_mm - 2 * margin
        available_height_mm = page_height_mm - 2 * margin

        # Calculate scaling factor to maintain aspect ratio
        width_ratio = available_width_mm / image_width_mm
        height_ratio = available_height_mm / image_height_mm
        scaling_factor = min(width_ratio, height_ratio)

        # Calculate the new image size in mm
        scaled_width_mm = image_width_mm * scaling_factor
        scaled_height_mm = image_height_mm * scaling_factor

        # Position the image with margins
        x_position = margin + (available_width_mm - scaled_width_mm) / 2
        y_position = margin + (available_height_mm - scaled_height_mm) / 2

        # Add image to PDF
        pdf.add_page()
        pdf.image(image, x=x_position, y=y_position, w=scaled_width_mm, h=scaled_height_mm)

    # Output PDF to bytes
    pdf_output = io.BytesIO()
    pdf.output(pdf_output)
    return pdf_output.getvalue()



def split_pdf_file(file: UploadFile):
    print(file.filename)
    pdf_reader = PdfReader(file.file)
    pdf_parts = []
    print(len(pdf_reader.pages))
    for i in range(len(pdf_reader.pages)):
        pdf_writer = PdfWriter()
        pdf_writer.add_page(pdf_reader.pages[i])

        pdf_page = io.BytesIO()
        pdf_writer.write(pdf_page)
        pdf_page.seek(0)
        pdf_parts.append((f"page_{i + 1}.pdf", pdf_page))

    return pdf_parts


def  pdf_to_ppt(pdf_bytes, keep_ratio, zoom_factor=2):
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    presentation = Presentation()
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        # Adjust zoom factor for higher resolution
        mat = fitz.Matrix(zoom_factor, zoom_factor)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")
        image_stream = io.BytesIO(img_bytes)
        image = Image.open(image_stream)
        image_buffer = io.BytesIO()
        image.save(image_buffer, format="PNG")
        image_buffer.seek(0)

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        if keep_ratio:
            image_width, image_height = image.size
            slide_ratio = slide_width / slide_height
            image_ratio = image_width / image_height

            if slide_ratio > image_ratio:
                height = slide_height
                width = int(height * image_ratio)
                left = int((slide_width - width) / 2)
                top = 0
            else:
                width = slide_width
                height = int(width / image_ratio)
                left = 0
                top = int((slide_height - height) / 2)

            slide.shapes.add_picture(image_buffer, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(image_buffer, 0, 0, width=slide_width, height=slide_height)

    ppt_buffer = io.BytesIO()
    presentation.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer


def pdf_to_word(pdf_bytes):
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    doc = Document()

    for page_num in range(len(pdf_document)):
        print(page_num)
        page = pdf_document.load_page(page_num)
        text = page.get_text()
        doc.add_paragraph(text)

    word_buffer = io.BytesIO()
    doc.save(word_buffer)
    word_buffer.seek(0)
    return word_buffer


def merge_pdfs(pdf_files):
    merged_pdf = fitz.open()

    for pdf_file in pdf_files:
        pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
        merged_pdf.insert_pdf(pdf_document)

    pdf_buffer = io.BytesIO()
    merged_pdf.save(pdf_buffer)
    pdf_buffer.seek(0)
    return pdf_buffer